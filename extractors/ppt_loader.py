from .base_file_loader import FileLoader
from pptx import Presentation

class PPTLoader(FileLoader):
    def validate_file(self, file_path):
        return file_path.endswith('.pptx')

    def load_file(self, file_path):
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return text, presentation.core_properties
