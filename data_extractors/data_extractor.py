from extractors.base_file_loader import FileLoader
from extractors.pdf_loader import PDFLoader
from extractors.docx_loader import DOCXLoader
from extractors.ppt_loader import PPTLoader
import docx
from pptx import Presentation

class DataExtractor:
    def __init__(self, loader: FileLoader):
        self.loader = loader
        self.text = None
        self.metadata = None

    def extract_text(self, file_path):
        self.text, self.metadata = self.loader.load_file(file_path)
        return self.text, self.metadata

    def extract_links(self, file_path):
        links = []
        if isinstance(self.loader, PDFLoader):
            # Example: Extract links from PDF
            for page in self.loader.load_file(file_path)[0]:
                # (Assuming you have some logic to extract hyperlinks)
                # Placeholder for actual link extraction logic
                links.append(("Sample Link Text", "http://example.com", 1))  # (text, url, page number)

        elif isinstance(self.loader, DOCXLoader):
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                for run in para.runs:
                    if run.hyperlink:  # Hypothetical check
                        links.append((run.text, run.hyperlink.target, para._element.getparent().index(para._element)))

        elif isinstance(self.loader, PPTLoader):
            presentation = Presentation(file_path)
            for slide_index, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "hyperlink") and shape.hyperlink:
                        links.append((shape.text, shape.hyperlink.address, slide_index))

        return links

    def extract_images(self, file_path):
        images = []
        if isinstance(self.loader, PDFLoader):
            # Extract images from PDF
            for page_index, page in enumerate(self.loader.load_file(file_path)[0]):
                # (Placeholder for actual image extraction logic)
                images.append(("image_path.png", "PNG", page_index))

        elif isinstance(self.loader, DOCXLoader):
            doc = docx.Document(file_path)
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    images.append((rel.target_ref, rel.target_ref.split('.')[-1].upper(), None))

        elif isinstance(self.loader, PPTLoader):
            presentation = Presentation(file_path)
            for slide_index, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        images.append((shape.image.filename, shape.image.ext, slide_index))

        return images

    def extract_tables(self, file_path):
        tables = []
        if isinstance(self.loader, PDFLoader):
            # Extract tables from PDF
            # (Placeholder for actual table extraction logic)
            tables.append(("Table Data", 1))  # (table data, page number)

        elif isinstance(self.loader, DOCXLoader):
            doc = docx.Document(file_path)
            for table_index, table in enumerate(doc.tables):
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                tables.append((table_data, table_index))

        elif isinstance(self.loader, PPTLoader):
            presentation = Presentation(file_path)
            for slide_index, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
                        tables.append((table_data, slide_index))

        return tables
